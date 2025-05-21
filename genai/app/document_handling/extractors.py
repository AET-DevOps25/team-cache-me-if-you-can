"""Enhanced text extraction utility (rawdict edition).

Uses PyMuPDF's *rawdict* mode to recover glyphs that regular
unicode extraction skips (e.g. mathematical symbols).  The rest of
the API and CLI interface are unchanged.
"""

from __future__ import annotations

import io
import logging
import os
from typing import IO
from collections import Counter

import fitz
from PyPDF2 import PdfReader
from pptx import Presentation
from docx import Document


logger = logging.getLogger(__name__)


class TextExtractor:
    """Extracts plain text from common document formats."""

    # ────────────────────────────── PDF ─────────────────────────────── #

    def _get_chars_from_rawdict_span(self, span: dict, font_name: str) -> list[str]:
        """Helper to extract characters from a rawdict span, handling CIDs."""
        chars_list: list[str] = []
        for ch in span.get("chars", []):
            c = ch.get("c")
            if c:  # direct Unicode
                chars_list.append(c)
            else:  # need CID‑to‑Unicode
                cid = ch.get("cid", 0)
                try:
                    uni = fitz.TOOLS.toUnicode(font_name, cid)
                    chars_list.append(uni or "")
                except Exception: 
                    chars_list.append("") 
        return chars_list

    def _extract_pdf_pymupdf(self, buffer: bytes) -> str:  # noqa: D401
        """
        Return plain text from *buffer* using PyMuPDF.
        Combines layout analysis ("blocks") with rawdict for character recovery within blocks,
        and attempts to filter common headers/footers.
        """
        all_blocks_info = [] 
        
        with fitz.open(stream=buffer, filetype="pdf") as doc:
            num_pages = doc.page_count
            for page_idx, page in enumerate(doc):
                page_rect = page.rect
                layout_blocks = page.get_text("blocks", sort=True)

                for lb_x0, lb_y0, lb_x1, lb_y1, block_text_simple, block_no, block_type in layout_blocks:
                    if block_type != 0: 
                        continue

                    block_bbox = fitz.Rect(lb_x0, lb_y0, lb_x1, lb_y1)
                    
                    raw_block_dict = page.get_text("rawdict", clip=block_bbox)
                    
                    block_chars_detailed: list[str] = []
                    for line in raw_block_dict.get("lines", []): 
                        for span in line.get("spans", []):
                            font_name = span["font"]
                            block_chars_detailed.extend(
                                self._get_chars_from_rawdict_span(span, font_name)
                            )

                    current_block_text = "".join(block_chars_detailed).strip()

                    if not current_block_text or len(current_block_text) < len(block_text_simple.strip()) / 2 :
                        current_block_text = block_text_simple.strip()
                    
                    if current_block_text:
                        all_blocks_info.append(
                            (page_idx, current_block_text, block_bbox, page_rect)
                        )
            
            common_hf_texts = set()
            if num_pages > 1: 
                block_text_counts = Counter(b_info[1] for b_info in all_blocks_info)
                
                min_occurrences_for_hf = max(2, int(num_pages * 0.3))
                
                for text, count in block_text_counts.items():
                    if count >= min_occurrences_for_hf and len(text) < 100:
                        in_hf_zone_count = 0
                        for p_idx, b_text, b_bbox, p_rect in all_blocks_info:
                            if b_text == text:
                                is_header_zone = b_bbox.y1 < p_rect.height * 0.15  # Top 15%
                                is_footer_zone = b_bbox.y0 > p_rect.height * 0.85  # Bottom 15%
                                if is_header_zone or is_footer_zone:
                                    in_hf_zone_count += 1
                        
                        # If >70% of its occurrences are in H/F zones, mark as H/F
                        if in_hf_zone_count / count > 0.7:
                            common_hf_texts.add(text)
                            logger.debug(f"Identified common H/F: '{text}'")
            
            # --- Reconstruct final text, skipping H/F and isolated page numbers ---
            output_by_page: list[list[str]] = [[] for _ in range(num_pages)]
            for p_idx, text, bbox, p_rect in all_blocks_info:
                # Check 1: Is it a common H/F text and in an H/F zone on this page?
                is_common_hf_in_zone = False
                if text in common_hf_texts:
                    is_header_zone = bbox.y1 < p_rect.height * 0.15
                    is_footer_zone = bbox.y0 > p_rect.height * 0.85
                    if is_header_zone or is_footer_zone:
                        is_common_hf_in_zone = True
                
                if is_common_hf_in_zone:
                    logger.debug(f"Skipping H/F block: '{text}' on page {p_idx}")
                    continue

                if text.strip().isdigit() and len(text.strip()) <= 4: # Max 4 digits for page num
                    is_extreme_top = bbox.y1 < p_rect.height * 0.08 # More stringent for page numbers
                    is_extreme_bottom = bbox.y0 > p_rect.height * 0.92
                    if is_extreme_top or is_extreme_bottom:
                        logger.debug(f"Skipping potential page number: '{text}' on page {p_idx}")
                        continue
                
                output_by_page[p_idx].append(text)

            page_strings = []
            for page_content_blocks in output_by_page:
                if page_content_blocks: 
                    page_strings.append("\n".join(page_content_blocks))
            
            return "\n\n".join(page_strings).strip()


    def _extract_pdf_pypdf2(self, file_io: IO[bytes]) -> str:
        """Return plain text using PyPDF2 (fallback)."""
        # This method might also benefit from the generic post-processor
        reader = PdfReader(file_io)
        text = "\n".join(
            filter(None, (page.extract_text() for page in reader.pages))
        ).strip()
        return text # Post-processing will be applied later by the dispatcher

    def extract_from_pdf(self, file_io: IO[bytes]) -> str:
        """Extract text from a PDF *file_io* stream."""
        file_io.seek(0)
        data = file_io.read()
        try:
            extracted_text = self._extract_pdf_pymupdf(data)
        except Exception as exc:
            logger.warning(
                "PyMuPDF rawdict extraction failed (%s). Falling back to PyPDF2.",
                exc,
                exc_info=True,
            )
            file_io.seek(0) # Reset for PyPDF2
            extracted_text = self._extract_pdf_pypdf2(file_io) # Pass the original file_io
        
        return extracted_text # Post-processing will be applied by the main dispatcher

    def _post_process_text(self, text: str) -> str:
        """Generic post-processing for extracted text."""
        if not text:
            return ""

        import re

        # 1. Normalize whitespace and basic line structure
        lines = text.splitlines()
        
        processed_lines = []
        for line in lines:
            line = line.strip()
            if line:
                processed_lines.append(line)
        
        paragraphs = re.split(r'\n\s*\n', text) 
        
        reconstructed_paragraphs = []
        for para_text in paragraphs:
            if not para_text.strip():
                continue

            current_para_lines = para_text.splitlines()
            if not current_para_lines:
                continue

            for i in range(len(current_para_lines) - 1):
                if current_para_lines[i].endswith('-'):
                    merged_line_start = current_para_lines[i][:-1] 
                    # Check if next line starts with a character that could continue a word
                    if current_para_lines[i+1] and current_para_lines[i+1][0].islower():
                        current_para_lines[i] = merged_line_start + current_para_lines[i+1]
                        current_para_lines[i+1] = "" # Mark for removal
            
            current_para_lines = [line for line in current_para_lines if line] # Remove emptied lines

            if not current_para_lines:
                reconstructed_paragraphs.append("") # Preserve paragraph break if original was just newlines
                continue

            new_paragraph_content = [current_para_lines[0]]
            for i in range(1, len(current_para_lines)):
                prev_line = new_paragraph_content[-1].strip()
                curr_line = current_para_lines[i].strip()

                if not curr_line: # Skip empty lines within a paragraph block
                    continue

                prev_ends_punct = prev_line.endswith(('.', '!', '?', ':', ';', ')', ']'))
                curr_starts_lower = curr_line and curr_line[0].islower()
                # More conservative joining:
                prev_ends_punct = prev_line.endswith(('.', '!', '?', ':', ';', ')', ']'))
                curr_starts_lower = curr_line and curr_line[0].islower()
                # Avoid joining list items or short headings to subsequent lines
                prev_is_list_item = prev_line.startswith(('-', '*', '•')) or re.match(r"^\s*\d+\.\s+", prev_line)
                # prev_is_short_heading = len(prev_line.split()) < 5 and prev_line.endswith(tuple("abcdefghijklmnopqrstuvwxyzABCDEFGHIJKLMNOPQRSTUVWXYZ0123456789"))

                if prev_line and not prev_ends_punct and curr_starts_lower and not prev_is_list_item:
                    new_paragraph_content[-1] = prev_line + " " + curr_line
                else:
                    new_paragraph_content.append(curr_line)
            
            reconstructed_paragraphs.append("\n".join(new_paragraph_content))

        # Join paragraphs with double newlines
        final_text = "\n\n".join(filter(None, reconstructed_paragraphs))
        
        # Final cleanup: replace multiple spaces with single space
        final_text = re.sub(r" +", " ", final_text)
        final_text = final_text.strip() # Remove leading/trailing whitespace/newlines from the whole text
        return final_text


    # ───────────────────────────── PPTX ─────────────────────────────── #

    def extract_from_pptx(self, file_io: IO[bytes]) -> str:
        """Extract concatenated text from all shapes in a PPTX file."""
        file_io.seek(0)
        prs = Presentation(file_io)
        texts: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if getattr(shape, "text", None):
                    texts.append(shape.text)
        return "\n".join(texts).strip()

    # ───────────────────────────── DOCX ─────────────────────────────── #

    def extract_from_docx(self, file_io: IO[bytes]) -> str:
        """Extract text from a modern Word (.docx) document."""
        file_io.seek(0)
        doc = Document(file_io)
        return "\n".join(p.text for p in doc.paragraphs if p.text).strip()

    # ────────────────────────── Dispatcher ──────────────────────────── #

    def extract_text(self, file_io: IO[bytes], filename: str) -> str:
        """Dispatch to the correct extractor based on *filename* extension."""
        ext = filename.lower().rsplit(".", 1)[-1]
        raw_text = ""
        if ext == "pdf":
            raw_text = self.extract_from_pdf(file_io)
        elif ext == "pptx":
            raw_text = self.extract_from_pptx(file_io)
        elif ext == "docx":
            raw_text = self.extract_from_docx(file_io)
        else:
            raise ValueError(f"Unsupported file type: {ext}")
        return self._post_process_text(raw_text)



if __name__ == "__main__":
    logging.basicConfig(level=logging.INFO) 

    extractor = TextExtractor()

    test_files = [
        "/Users/i_gore/PycharmProjects/study_sync/data/docs/mdl_intro.pdf",
    ]

    processed_dir = "/Users/i_gore/PycharmProjects/study_sync/data/processed"
    # Ensure processed_dir exists
    os.makedirs(processed_dir, exist_ok=True)

    for full_file_path in test_files:
        if not os.path.exists(full_file_path):
            print(f"SKIPPING: {full_file_path} not found.")
            continue
            
        try:
            file_name_with_ext = os.path.basename(full_file_path)
            print(f"\n--- Processing {file_name_with_ext} ---")
            with open(full_file_path, "rb") as f:
                extracted = extractor.extract_text(f, file_name_with_ext) # Pass filename for extension detection

            base_name_no_ext = os.path.splitext(file_name_with_ext)[0]
            output_filename = base_name_no_ext + ".txt"
            output_path = os.path.join(processed_dir, output_filename)
            
            with open(output_path, "w", encoding="utf-8") as out_file:
                out_file.write(extracted)

            print(f"--- {file_name_with_ext} processed successfully ---")
            print(f"Text saved to: {output_path}")


        except FileNotFoundError: # Should be caught by os.path.exists now
            print(f"Error: {full_file_path} not found for testing.")
        except Exception as err:
            print(f"Error processing {full_file_path}: {err}")
            logger.exception(f"Detailed error for {full_file_path}:")